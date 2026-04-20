# Here is a **line‑by‑line explanation** of the PowerPoint text extraction code, presented as bullet points:

# - `from pptx import Presentation`
#   - Imports the `Presentation` class from the `python-pptx` library.
#   - This class is used to open and manipulate PowerPoint (.pptx) files.

# - `ppt_path = "C:\\Personal\\2024\\Learning\\Generative AI\\RAG\\27_Context_Engineering\\2_RAG\\1_Document_Processing\\Data\\Prompt_Engineering_Guide.pptx"`
#   - Assigns a string containing the file path to a PowerPoint presentation named `Prompt_Engineering_Guide.pptx`.
#   - Double backslashes (`\\`) are used to escape backslashes in a regular string (or raw strings could have been used).

# - `def extract_ppt_text(ppt_path):`
#   - Defines a function named `extract_ppt_text` that takes one parameter `ppt_path` (the path to the PowerPoint file).
#   - The function will extract all text content from the presentation.

# - `prs = Presentation(ppt_path)`
#   - Creates a `Presentation` object by opening the PowerPoint file located at `ppt_path`.
#   - This object provides access to slides, shapes, and other elements of the presentation.

# - `all_text = ""`
#   - Initializes an empty string that will accumulate the extracted text from all slides.

# - `for i, slide in enumerate(prs.slides):`
#   - Loops over each slide in the presentation. `prs.slides` returns a list of slide objects.
#   - `enumerate` provides the slide index `i` (starting from 0) and the `slide` object.

# - `all_text += f"\n--- Slide {i+1} ---\n"`
#   - Appends a header line to `all_text` indicating the slide number (adds 1 to `i` for human‑readable numbering).
#   - The header is surrounded by newlines for separation.

# - `for shape in slide.shapes:`
#   - Loops over all shapes on the current slide. Shapes can include text boxes, images, charts, etc.

# - `if shape.has_text_frame:`
#   - Checks whether the current shape contains a text frame (i.e., can hold text). Some shapes (like images) do not have text.

# - `for paragraph in shape.text_frame.paragraphs:`
#   - Loops over each paragraph within the shape’s text frame. A paragraph is a block of text with optional formatting.

# - `all_text += paragraph.text + "\n"`
#   - Appends the text content of the current paragraph to `all_text`, followed by a newline character.
#   - `paragraph.text` returns a string of the plain text in that paragraph (without formatting).

# - `return all_text`
#   - After processing all slides, shapes, and paragraphs, returns the complete concatenated text string.

# - `text = extract_ppt_text(ppt_path)`
#   - Calls the `extract_ppt_text` function with the previously defined `ppt_path` and stores the returned text in the variable `text`.

# - `print(text)`
#   - Prints the entire extracted text (all slides with headers) to the console.

from pptx import Presentation

ppt_path = "D:\GenAI Content\\AI code\\4_RAG_Indexing\\1_Document_Processing\\Data\\Prompt_Engineering_Guide.pptx"

def extract_ppt_text(ppt_path):
    prs = Presentation(ppt_path)
    all_text = ""

    for i, slide in enumerate(prs.slides):
        all_text += f"\n--- Slide {i+1} ---\n"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    all_text += paragraph.text + "\n"

    return all_text


text = extract_ppt_text(ppt_path)
print(text)